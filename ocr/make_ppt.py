from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def set_slide_title(slide, title_text):
    """슬라이드에 제목을 추가합니다."""
    try:
        title = slide.shapes.title
        title.text = title_text
        title.text_frame.paragraphs[0].font.size = Pt(32)
        title.text_frame.paragraphs[0].font.bold = True
    except AttributeError:
        # 제목 슬라이드 레이아웃이 아닌 경우
        pass

def add_content_box(slide, content_list, left=Inches(0.5), top=Inches(1.5), width=Inches(9.0), height=Inches(5.5)):
    """슬라이드에 텍스트 콘텐츠 박스를 추가합니다."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for item in content_list:
        if item.startswith("• "):  # 기본 글머리 기호
            p = tf.add_paragraph()
            p.text = item[2:]
            p.level = 0
            p.font.size = Pt(20)
            p.font.name = 'Malgun Gothic'
        elif item.startswith("  - "): # 하위 글머리 기호
            p = tf.add_paragraph()
            p.text = item[4:]
            p.level = 1
            p.font.size = Pt(18)
            p.font.name = 'Malgun Gothic'
        else: # 소제목 또는 일반 텍스트
            p = tf.add_paragraph()
            p.text = item
            p.font.size = Pt(22)
            p.font.bold = True
            p.font.name = 'Malgun Gothic'
            p.space_after = Pt(12)

    return tf

def create_presentation():
    """OCR 비교 분석 PPT를 생성합니다."""
    prs = Presentation()
    
    # 슬라이드 레이아웃 정의
    title_slide_layout = prs.slide_layouts[0] # 제목 슬라이드
    content_layout = prs.slide_layouts[1]     # 제목 + 콘텐츠
    blank_layout = prs.slide_layouts[6]        # 빈 슬라이드

    # --- 슬라이드 1: 표지 (지침 4.1 일부) ---
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "OCR 기술별 한글 인식률 비교 분석"
    subtitle.text = "Tesseract, OCR.space API, PyMuPDF 활용\n발표자: 이지우"

    # --- 슬라이드 2: OCR 개요 (정의 및 원리) (지침 4.1) ---
    slide = prs.slides.add_slide(content_layout)
    set_slide_title(slide, "1. OCR 개요 (정의 및 원리)")
    content = [
        "OCR (Optical Character Recognition, 광학 문자 인식)",
        "• 이미지(사진, 스캔) 속 문자를 기계가 읽을 수 있는 텍IST로 변환하는 기술",
        "작동 원리",
        "• 이미지 전처리: 노이즈 제거, 기울기 보정",
        "• 문자 탐지 (Layout Analysis): 텍스트 영역 탐색",
        "• 문자 인식 (Recognition): 딥러닝/머신러닝 모델로 문자 판독",
        "• 후처리 (Post-processing): 오류 수정 및 텍스트 조합",
    ]
    add_content_box(slide, content)

    # --- 슬라이드 3: OCR 활용 분야 (지침 4.2) ---
    slide = prs.slides.add_slide(content_layout)
    set_slide_title(slide, "2. OCR 활용 분야 (사례)")
    content = [
        "금융",
        "• 신분증/운전면허증 촬영 본인 인증, 신용카드 번호 인식",
        "교육",
        "• 교재/논문 스캔 및 검색, 필기노트 디지털화, OMR 채점",
        "행정 (공공)",
        "• 민원 서류(등본, 증명서) 자동 처리, 우편물 주소 인식 및 분류",
        "의료 / 법률 / 물류",
        "• 의료: 처방전, 진료기록부 텍스트 변환",
        "• 법률: 과거 판례, 법률 문서 검색 데이터베이스 구축",
        "• 물류: 택배 송장 번호 인식, 차량 번호판 자동 인식",
    ]
    add_content_box(slide, content)

    # --- 슬라이드 4: OCR 처리 과정 설명 (지침 4.3) ---
    slide = prs.slides.add_slide(content_layout)
    set_slide_title(slide, "3. 프로젝트 OCR 처리 과정")
    content = [
        "방식 1: Tesseract (로컬 엔진)",
        "• `pytesseract` 라이브러리 활용, `01.jpg` 파일 직접 처리",
        "• Mac Homebrew로 `tesseract-lang` (한국어) 설치",
        "방식 2: OCR.space API (웹 서비스)",
        "• `requests` 라이브러리 활용, `01.jpg` 파일을 API 서버로 전송/수신",
        "방식 3: PyMuPDF + Tesseract (PDF 처리)",
        "• `fitz`로 `sample.pdf` 파일 열기",
        "• 텍스트 레이어는 `page.get_text()`로 직접 추출",
        "• PDF 내 이미지는 `page.get_images()`로 분리 후 Tesseract로 OCR 수행",
    ]
    add_content_box(slide, content)

    # --- 슬라이드 5: 인식률 평가 방법 (지침 3.1, 4.4) ---
    slide = prs.slides.add_slide(content_layout)
    set_slide_title(slide, "4. 인식률 평가 방법 (Levenshtein)")
    content = [
        "Levenshtein 거리 (편집 거리)",
        "• '정답' 텍스트와 'OCR 결과' 텍스트가 얼마나 다른지 측정",
        "• 한 문자열을 다른 문자열로 바꾸는 데 필요한 최소 수정 횟수 (삽입, 삭제, 대체)",
        "• 예: 'OCR' → 'OCF' (거리: 1)",
        "인식률(Accuracy) 계산 공식",
        "• `정확도(%) = (1 - (거리 / 더 긴 텍스트 길이)) * 100`",
        "• 거리가 0이면 정확도 100%",
        "Python 예시 코드",
        "  - `distance = Levenshtein.distance(original, ocr_result)`",
        "  - `max_len = max(len(original), len(ocr_result))`",
        "  - `accuracy = (1 - distance / max_len) * 100`",
    ]
    add_content_box(slide, content)

    # --- 슬라이드 6: 인식률 결과 그래프 (지침 1, 3.2, 4.5) ---
    slide = prs.slides.add_slide(content_layout)
    set_slide_title(slide, "5. 인식률 결과 및 원본 비교 그래프")
    
    # 이미지 파일 경로
    img_path = 'ocr/ocr_comparison_chart.png'

    try:
        # 그래프 이미지 삽입 (가운데 정렬)
        pic = slide.shapes.add_picture(img_path, Inches(0.5), Inches(1.5), width=Inches(9.0))
    except FileNotFoundError:
        add_content_box(slide, ["오류: 'ocr_comparison_chart.png' 파일을 찾을 수 없습니다."])
    except Exception as e:
         add_content_box(slide, [f"이미지 로드 중 오류 발생: {e}"])

    # (결과 텍스트를 이곳에 추가할 수 있습니다)
    # 예: add_content_box(slide, ["Pytesseract(JPG): XX.X%", ...], top=Inches(5.0))


    # --- 슬라이드 7: 결론 및 향후 발전 방향 (지침 4.6) ---
    slide = prs.slides.add_slide(content_layout)
    set_slide_title(slide, "6. 결론 및 향후 발전 방향")
    content = [
        "결론",
        "• (결과에 따라 작성: 예: 'OCR.space API'가 JPG 이미지에서 가장 높은 한글 인식률을 보임)",
        "• (결과에 따라 작성: 예: 'Tesseract'는 설치형임에도 준수한 성능을 제공함)",
        "• 'PyMuPDF'는 텍스트 레이어 추출이 가능할 때 가장 빠르고 정확함",
        "한계점",
        "• 제한된 데이터셋(JPG 1장, PDF 1장)으로 테스트됨",
        "• 손글씨, 저화질, 복잡한 배경 등 다양한 변수 미적용",
        "향후 발전 방향",
        "• Tesseract 모델 파인튜닝(Fine-tuning)을 통한 특정 폰트 인식률 개선",
        "• 다양한 이미지 전처리(Pre-processing) 기법 적용 테스트",
    ]
    add_content_box(slide, content)

    # --- 슬라이드 8: Q&A ---
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "감사합니다"
    subtitle.text = "Q & A"

    # --- 파일 저장 ---
    file_path = "OCR_비교_분석_발표자료.pptx"
    prs.save(file_path)
    return file_path

if __name__ == '__main__':
    try:
        output_file = create_presentation()
        print(f"성공! '{output_file}' 파일이 생성되었습니다.")
    except Exception as e:
        print(f"PPT 생성 중 오류가 발생했습니다: {e}")
        if "ocr_comparison_chart.png" in str(e):
            print("오류: 'ocr_comparison_chart.png' 파일이 스크립트와 같은 폴더에 있는지 확인하세요.")
        if "python-pptx" in str(e):
             print("오류: 'python-pptx' 라이브러리가 설치되었는지 확인하세요. (pip install python-pptx)")